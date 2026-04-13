from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_presentation():
    prs = Presentation()

    # Define color scheme (Dark theme like friend's PPT)
    bg_color = RGBColor(10, 25, 41) # Dark Blue
    text_color = RGBColor(255, 255, 255) # White
    accent_color = RGBColor(123, 47, 255) # Purple accent

    def set_slide_background(slide):
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color

    def add_page_number(slide, page_num):
        txBox = slide.shapes.add_textbox(Inches(9), Inches(0.2), Inches(1), Inches(0.5))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = str(page_num)
        p.font.size = Pt(12)
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.RIGHT

    def add_enrollment_footer(slide):
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(7.1), Inches(3), Inches(0.3))
        tf = txBox.text_frame
        p = tf.paragraphs[0]
        p.text = "245300694048"
        p.font.size = Pt(10)
        p.font.color.rgb = RGBColor(150, 150, 150)

    # --- Slide 1: Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    
    # Title
    txBox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(2))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "IntelliScan"
    p.font.bold = True
    p.font.size = Pt(64)
    p.font.color.rgb = text_color
    
    # Subtitle
    p = tf.add_paragraph()
    p.text = "AI-Powered Business Card Management & CRM Automation"
    p.font.size = Pt(24)
    p.font.color.rgb = accent_color

    # Info

    
    txBox = slide.shapes.add_textbox(Inches(6), Inches(5.5), Inches(4), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = "Guide: PROF. Nil Gosai"
    p.font.size = Pt(18)
    p.font.color.rgb = text_color

    # --- Slide 2: Agenda ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 2)
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(4), Inches(1))
    title.text_frame.text = "Agenda"
    title.text_frame.paragraphs[0].font.size = Pt(44)
    title.text_frame.paragraphs[0].font.color.rgb = text_color

    agenda_items = [
        "01  Introduction and Industry Context",
        "02  Overview of the Platform",
        "03  Core Platform Features",
        "04  Workflow of the System",
        "05  Tools & Technology Used",
        "06  System Design",
        "07  Development Process",
        "08  Agile Documentation",
        "09  Development Screenshots",
        "10  Proposed Enhancements"
    ]
    
    y = 1.8
    for item in agenda_items:
        tx = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.5))
        p = tx.text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(20)
        p.font.color.rgb = text_color
        y += 0.5

    # --- Slide 3: Introduction ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 3)
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    title.text_frame.text = "Introduction & Industry Context"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    title.text_frame.paragraphs[0].font.color.rgb = text_color
    
    content = [
        "• Rapid growth of AI-driven lead capture technologies.",
        "• Fragmented physical-to-digital bridge in B2B sales.",
        "• Need for frictionless, 'Zero-Touch' data entry.",
        "• Transitioning from manual logging to automated intelligent workflows."
    ]
    y = 2.0
    for line in content:
        tx = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.5))
        p = tx.text_frame.paragraphs[0]
        p.text = line
        p.font.size = Pt(22)
        p.font.color.rgb = text_color
        y += 0.8

    # --- Slide 4: About IntelliScan ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 4)
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    title.text_frame.text = "About IntelliScan"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(8), Inches(0.5))
    subtitle.text_frame.text = "Overview of the Platform"
    subtitle.text_frame.paragraphs[0].font.size = Pt(20)
    subtitle.text_frame.paragraphs[0].font.color.rgb = accent_color
    
    body = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
    tf = body.text_frame
    p = tf.paragraphs[0]
    p.text = "IntelliScan is an AI-powered Business Card Scanner and intelligent CRM automation platform that centralizes data extraction and eliminates manual data entry."
    p.font.size = Pt(22)
    p.font.color.rgb = text_color
    
    p = tf.add_paragraph()
    p.text = "\nThe platform supports:\n• Scene management and bulk lead uploads\n• Real-time AI extraction (Gemini Vision)\n• Automated pipeline routing\n• Engagement telemetry"
    p.font.size = Pt(20)

    # --- Slide 5: Core Features ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 5)
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    title.text_frame.text = "Core Platform Features"
    title.text_frame.paragraphs[0].font.size = Pt(36)

    features = [
        ("AI Card Scanning", "OCR and Google Gemini integration for attribute extraction."),
        ("Lead Routing Engine", "Dynamic assignment based on inferred seniority."),
        ("Email Marketing Hub", "AI-authored email drafting and drip sequences."),
        ("Data Quality Center", "Fuzzy matching and deduplication alerts."),
        ("Public Booking Tool", "Self-service calendar scheduling via digital cards.")
    ]
    y = 2.0
    for f_title, f_desc in features:
        tx = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.5))
        p = tx.text_frame.paragraphs[0]
        p.text = f"{f_title}: {f_desc}"
        p.font.size = Pt(18)
        p.font.color.rgb = text_color
        y += 0.7

    # --- Slide 6: End-to-End Workflow ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 6)
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    title.text_frame.text = "End-to-End Authoring Workflow"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    steps = ["Login & Dashboard", "Capture Card Image", "AI Metadata Inference", "CRM Pipeline Entry", "Automated Follow-up"]
    y = 2.5
    for i, step in enumerate(steps):
        tx = slide.shapes.add_textbox(Inches(1 + (i*1.5)), Inches(y), Inches(2), Inches(1))
        p = tx.text_frame.paragraphs[0]
        p.text = f"Step {i+1}\n{step}"
        p.font.size = Pt(14)
        p.font.color.rgb = text_color
        p.alignment = PP_ALIGN.CENTER

    # --- Slide 7: Tools & Technology ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 7)
    
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    title.text_frame.text = "Tools & Technologies Used"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    techs = {
        "Frontend": "React, Vite, Tailwind CSS, Lucide",
        "Backend": "Node.js, Express JS, REST APIs",
        "Database": "SQLite, Local Persistence",
        "AI Engine": "Google Gemini Pro Vision LLM"
    }
    y = 2.5
    for cat, val in techs.items():
        tx = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.5))
        p = tx.text_frame.paragraphs[0]
        p.text = f"{cat}: {val}"
        p.font.size = Pt(22)
        p.font.color.rgb = text_color
        y += 0.8

    # --- Slide 8: System Design ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 8)
    title = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
    p = title.text_frame.paragraphs[0]
    p.text = "System Design & UML Diagrams"
    p.font.size = Pt(54)
    p.alignment = PP_ALIGN.CENTER

    # --- Slides 9-14: UML Placeholders ---
    diagrams = ["Use Case Diagram (Global)", "User Authentication Flow", "AI Vision Pipeline", "Sequence Diagram (API Routing)", "Class Diagram", "Database Architecture"]
    for i, diag in enumerate(diagrams):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide)
        add_page_number(slide, 9 + i)
        add_enrollment_footer(slide)
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
        title.text_frame.text = f"{i+1}. {diag}"
        title.text_frame.paragraphs[0].font.size = Pt(28)
        
        # Placeholder for image
        tx = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
        p = tx.text_frame.paragraphs[0]
        p.text = "[PASTE DIAGRAM HERE FROM REPORT]"
        p.font.size = Pt(24)
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.CENTER

    # --- Slides 15-20: Data Dictionary ---
    tables = [
        ("Workspaces", "Core tenant metadata including IDs and quotas."),
        ("Users", "Authentication matrix and RBAC role definitions."),
        ("Contacts", "Extracted lead data mapped from OCR processes."),
        ("Scanned_Images", "Raw storage matrix for upload/inference logs."),
        ("Events", "Trade-show and networking function groupings."),
        ("Email_Campaigns", "Background queue nodes for SMTP relays.")
    ]
    for i, (table, desc) in enumerate(tables):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide)
        add_page_number(slide, 15 + i)
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
        title.text_frame.text = f"Data Dictionary: {table} Table"
        title.text_frame.paragraphs[0].font.size = Pt(28)
        
        tx = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1))
        p = tx.text_frame.paragraphs[0]
        p.text = desc
        p.font.size = Pt(20)
        p.font.color.rgb = text_color

    # --- Slide 21: Development ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 21)
    title = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
    p = title.text_frame.paragraphs[0]
    p.text = "Development & UI Screenshots"
    p.font.size = Pt(54)
    p.alignment = PP_ALIGN.CENTER

    # --- Slides 22-29: Screenshot Placeholders ---
    pages = ["Login Screen", "Main Dashboard", "Scanning Interface", "AI Coach Insights", "Lead Pipeline (Kanban)", "Email Marketing Hub", "Settings & RBAC", "Public Digital Card"]
    for i, page in enumerate(pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide)
        add_page_number(slide, 22 + i)
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
        title.text_frame.text = f"Platform UI: {page}"
        title.text_frame.paragraphs[0].font.size = Pt(28)
        
        # Placeholder for image
        tx = slide.shapes.add_textbox(Inches(1), Inches(3.5), Inches(8), Inches(1))
        p = tx.text_frame.paragraphs[0]
        p.text = "[PASTE SCREENSHOT HERE]"
        p.font.size = Pt(24)
        p.font.color.rgb = accent_color
        p.alignment = PP_ALIGN.CENTER

    # --- Slide 30: Agile Intro ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 30)
    title = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(2))
    p = title.text_frame.paragraphs[0]
    p.text = "Agile Methodology"
    p.font.size = Pt(54)
    p.alignment = PP_ALIGN.CENTER

    # --- Slide 31: Agile Project Charter ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 31)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    title.text_frame.text = "Agile Project Charter"
    title.text_frame.paragraphs[0].font.size = Pt(36)
    
    charter = [
        "• Project Scope: AI-powered B2B lead management.",
        "• Success Criteria: Zero-touch extraction & automated drips.",
        "• Risks: API latency, data complexity.",
        "• Mitigation: Modular architecture & fuzzy-logic matching."
    ]
    y = 2.0
    for line in charter:
        tx = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(0.5))
        p = tx.text_frame.paragraphs[0]
        p.text = line
        p.font.size = Pt(20)
        y += 0.8

    # --- Slide 32-33: Project Timeline (Q1-Q8) ---
    quarters = [
        "Q1-Q2: Planning & Backend Setup (Authentication, DB Schema)",
        "Q3-Q4: AI Vision Integration & CRM Pipeline Building",
        "Q5-Q6: Email Hub Development & Theme Rollout",
        "Q7-Q8: Final Testing, Documentation, and Submission"
    ]
    for i in range(2):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        set_slide_background(slide)
        add_page_number(slide, 32 + i)
        title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
        title.text_frame.text = f"Project Timeline (Part {i+1})"
        y = 2.0
        idx = i * 2
        for line in quarters[idx:idx+2]:
            tx = slide.shapes.add_textbox(Inches(1), Inches(y), Inches(8), Inches(1))
            tx.text_frame.text = line
            y += 1.5

    # --- Slide 34: Project Plan ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 34)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    title.text_frame.text = "Agile Project Plan Overview"
    tx = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tx.text_frame.text = "• Planned Start: Jan 2026\n• Planned End: April 2026\n• Methodology: Iterative Sprint Sprints (S1-S8)\n• Metrics: Bug resolution rate, OCR accuracy."

    # --- Slide 35: Risks ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 35)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    title.text_frame.text = "Success Criteria & Risks"
    tx = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tx.text_frame.text = "Risks:\n- Handling diverse low-quality images.\n- UI responsiveness on large datasets.\nMitigation:\n- Implementation of high-precision LLM inference.\n- Fragmented state management via Context API."

    # --- Slide 36: User Stories ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 36)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    title.text_frame.text = "Agile User Stories"
    tx = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tx.text_frame.text = "1. AI Scanning: Develop dynamic card extraction.\n2. CRM Pipeline: Advanced repository with smart search.\n3. Sequences: AI email automation engine.\n4. Kiosk Module: Bulk lead capture in networking mode."

    # --- Slide 37-38: Backlog ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 37)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    title.text_frame.text = "Agile Sprint Backlog (S1-S4)"
    tx = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tx.text_frame.text = "S1: Setup & Foundational APIs\nS2: Core Dashboard Implementation\nS3: AI Vision Integration\nS4: Pipeline & Filtering Logic"

    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 38)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    title.text_frame.text = "Agile Sprint Backlog (S5-S8)"
    tx = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tx.text_frame.text = "S5: Automation & Daemons\nS6: Calendar & Kiosk Integration\nS7: Multi-tenant Billing Config\nS8: Documentation & Phase-2 Finalization"

    # --- Slide 39: Earned Values ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 39)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    title.text_frame.text = "Earned Values & Burned Chart"
    tx = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(8), Inches(3))
    tx.text_frame.text = "Planned Progress vs Actual Delivery:\nSprint 1-4: High performance (Value 14.0 peaking at S4).\nSprint 5-8: Sustained maintenance value."

    # --- Slide 40: Future Enhancements ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 40)
    title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(8), Inches(1))
    title.text_frame.text = "Future Enhancements"
    tx = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(3))
    tx.text_frame.text = "• Native iOS/Android Mobile Application\n• LinkedIn Automation Integration\n• Real-time Prospect Alerts via Push Notifs\n• Multi-language Expansion"

    # --- Slide 41: Thank You ---
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide)
    add_page_number(slide, 41)
    
    msg = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(8), Inches(2))
    p = msg.text_frame.paragraphs[0]
    p.text = "Thank You!"
    p.font.size = Pt(80)
    p.font.bold = True
    p.alignment = PP_ALIGN.CENTER
    
    info = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(8), Inches(1))
    p = info.text_frame.paragraphs[0]
    p.text = "- ANANT PRABHUDESAI"
    p.font.size = Pt(24)
    p.alignment = PP_ALIGN.CENTER

    prs.save('IntelliScan_Final_Presentation.pptx')
    print("Presentation created successfully!")

if __name__ == "__main__":
    create_presentation()
