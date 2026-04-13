#!/usr/bin/env python3
"""
IntelliScan Presentation Generator
Creates a PowerPoint presentation from the IntelliScan documentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

def create_title_slide(prs, title, subtitle):
    """Create a title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_box = slide.shapes.title
    subtitle_box = slide.placeholders[1]

    title_box.text = title
    subtitle_box.text = subtitle

    # Style the title
    title_box.text_frame.paragraphs[0].font.size = Pt(44)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    return slide

def create_content_slide(prs, title, content_list):
    """Create a content slide with bullet points"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_box = slide.shapes.title
    content_box = slide.placeholders[1]

    title_box.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True

    for item in content_list:
        p = content_box.text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(24)
        p.level = 0

    return slide

def create_two_column_slide(prs, title, left_title, left_content, right_title, right_content):
    """Create a slide with two columns"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_box = slide.shapes.title
    title_box.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    title_box.text_frame.paragraphs[0].font.bold = True

    # Create two text boxes
    left_box = slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(4.5), Inches(5))
    right_box = slide.shapes.add_textbox(Inches(5.5), Inches(1.5), Inches(4.5), Inches(5))

    # Left column title
    left_title_para = left_box.text_frame.add_paragraph()
    left_title_para.text = left_title
    left_title_para.font.bold = True
    left_title_para.font.size = Pt(24)

    for item in left_content:
        p = left_box.text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)

    # Right column title
    right_title_para = right_box.text_frame.add_paragraph()
    right_title_para.text = right_title
    right_title_para.font.bold = True
    right_title_para.font.size = Pt(24)

    for item in right_content:
        p = right_box.text_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)

    return slide

def create_intelliscan_presentation():
    """Create the complete IntelliScan presentation"""

    # Create presentation
    prs = Presentation()

    # Slide 1: Title Slide
    create_title_slide(
        prs,
        "IntelliScan: AI-Powered Business Card Scanning & CRM Platform",
        "Major Project Presentation\nMCA Semester 4\nApril 4, 2026"
    )

    # Slide 2: Team Details
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_box = slide.shapes.title
    content_box = slide.placeholders[1]

    title_box.text = "Project Team"
    title_box.text_frame.paragraphs[0].font.size = Pt(32)

    team_info = [
        "Team Members: [Fill Names]",
        "Enrollment Numbers: [Fill Enrollment Nos.]",
        "Guide: Prof. Khushbu Patel",
        "Presentation Date: April 4, 2026"
    ]

    for item in team_info:
        p = content_box.text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(24)

    # Slide 3: Agenda
    agenda_items = [
        "1. Introduction & Existing System Analysis",
        "2. Need for the New System",
        "3. Problem Definition & Objectives",
        "4. System Scope & Features",
        "5. Technology Stack",
        "6. System Architecture",
        "7. Database Design",
        "8. Implementation & Demo",
        "9. Testing & Future Scope"
    ]
    create_content_slide(prs, "Presentation Agenda", agenda_items)

    # Slide 4: Existing System Problems
    problems = [
        "Manual data entry into spreadsheets/CRMs (slow & error-prone)",
        "Physical storage of business cards (risk of loss, no search capability)",
        "Fragmented tools - contacts, emails, notes, events not connected",
        "Weak follow-up hygiene - leads go stale due to lack of automation",
        "No enterprise collaboration features for team-based CRM"
    ]
    create_content_slide(prs, "Existing System Analysis", problems)

    # Slide 5: Need for New System
    needs = [
        "Instant digitization - AI OCR converts cards to structured data",
        "Centralized CRM with search, edit, export, and event tagging",
        "Enterprise workflows - workspace collaboration, data deduplication",
        "Automated follow-up - AI drafts, email marketing, calendar booking",
        "Tiered billing system (Razorpay) for different user levels"
    ]
    create_content_slide(prs, "Need for the New System", needs)

    # Slide 6: Problem Definition
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_box = slide.shapes.title
    content_box = slide.placeholders[1]

    title_box.text = "Problem Definition & Objectives"
    title_box.text_frame.paragraphs[0].font.size = Pt(32)

    problem_text = [
        "Problem: Networking leakage occurs when valuable business leads are lost because:",
        "• Business card capture is slow and error-prone",
        "• Physical cards are difficult to organize and search",
        "• Follow-up processes are inconsistent and manual",
        "• Enterprise teams lack collaborative CRM tools"
    ]

    objectives = [
        "",
        "Objectives:",
        "• Build AI-powered single and multi-card scanners",
        "• Create centralized CRM with role-based access control",
        "• Implement automated follow-up workflows",
        "• Provide enterprise data governance tools"
    ]

    for item in problem_text + objectives:
        p = content_box.text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(20)

    # Slide 7: System Scope - Core Features
    features = [
        "✓ Authentication & Session Management",
        "✓ Single Card AI Scan (Gemini → OpenAI → Tesseract fallback)",
        "✓ Group Photo Multi-Card Scan (Enterprise)",
        "✓ Contacts CRM (Search, Edit, Export, Deal Scoring)",
        "✓ Events & Campaigns Management",
        "✓ AI Drafts & Networking Coach",
        "✓ Email Marketing (Templates, Lists, Campaigns)",
        "✓ Calendar & Booking Links",
        "✓ Data Quality Center (Deduplication/Merge)",
        "✓ Compliance Policies & Audit Logs",
        "✓ Third-party Integrations (Webhooks, CRM Mapping)",
        "✓ Billing & Plan Upgrades (Razorpay)"
    ]
    create_content_slide(prs, "System Scope - Implemented Features", features)

    # Slide 8: Technology Stack
    frontend_stack = [
        "React 18 (Vite build tool)",
        "Tailwind CSS (Utility-first styling)",
        "React Router (Client-side routing)",
        "Axios (HTTP client)",
        "Lucide React (Icons)",
        "Context API (State management)"
    ]

    backend_stack = [
        "Node.js + Express.js",
        "SQLite3 (Database)",
        "Socket.IO (Real-time features)",
        "JWT (Authentication)",
        "bcryptjs (Password hashing)"
    ]

    create_two_column_slide(
        prs,
        "Technology Stack",
        "Frontend (React SPA)",
        frontend_stack,
        "Backend (Express API)",
        backend_stack
    )

    # Slide 9: AI & External Services
    ai_services = [
        "Google Gemini 1.5/2.0 Flash (Primary AI vision)",
        "OpenAI GPT (Fallback text/vision processing)",
        "Tesseract.js (Offline OCR fallback)",
        "Razorpay (Payment processing)",
        "Nodemailer (SMTP email sending)"
    ]
    create_content_slide(prs, "AI & Third-Party Integrations", ai_services)

    # Slide 10: System Architecture
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_box = slide.shapes.title
    content_box = slide.placeholders[1]

    title_box.text = "System Architecture"
    title_box.text_frame.paragraphs[0].font.size = Pt(32)

    arch_text = [
        "Frontend (intelliscan-app):",
        "• React SPA with Vite build system",
        "• Component-based architecture with layouts",
        "• Context providers for auth, contacts, batch queue",
        "• Responsive design with Tailwind CSS",
        "",
        "Backend (intelliscan-server):",
        "• Express.js REST API with middleware",
        "• SQLite relational database",
        "• JWT authentication with role-based access",
        "• Socket.IO for real-time collaboration",
        "• Unified AI extraction pipeline with fallbacks"
    ]

    for item in arch_text:
        p = content_box.text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(20)

    # Slide 11: User Roles & Tiers
    roles = [
        "Anonymous/Public: Landing page, pricing, public analytics",
        "Personal User (Free/Pro): Single scans, basic CRM, AI drafts",
        "Enterprise User: Multi-card scans, email marketing, calendar",
        "Enterprise Admin: Data policies, integrations, team management",
        "Super Admin: Platform monitoring, model tuning, incident management"
    ]
    create_content_slide(prs, "User Roles & Access Control", roles)

    # Slide 12: Database Design
    db_tables = [
        "users: id, email, password, role, tier, created_at",
        "user_quotas: user_id, overall_limit, used_count, group_scans_limit, group_scans_used",
        "contacts: id, user_id, json_data, deal_score, scan_date, event_id",
        "events: id, user_id, name, description, created_at",
        "ai_drafts: id, contact_id, draft_type, content, created_at",
        "email_campaigns: id, user_id, name, template_id, status",
        "audit_logs: id, user_id, action, resource_type, timestamp"
    ]
    create_content_slide(prs, "Database Schema (SQLite)", db_tables)

    # Slide 13: Core Workflows
    workflows = [
        "Single Card Scan: Upload → Validate → AI Extract → Preview → Save",
        "Group Photo Scan: Upload → Enterprise Check → Multi-Extract → Batch Save",
        "Contact Management: Search → Edit → Tag Events → Export → Enrich",
        "AI Follow-up: Select Contact → Generate Draft → Customize → Send",
        "Email Marketing: Create List → Design Template → Schedule Campaign",
        "Billing Upgrade: Select Plan → Razorpay Order → Payment → Unlock Features"
    ]
    create_content_slide(prs, "Core Business Workflows", workflows)

    # Slide 14: Demo Plan
    demo_steps = [
        "1. Landing Page & Authentication",
        "2. Single Card Scan Demo",
        "3. Contacts CRM Overview",
        "4. AI Drafts Generation",
        "5. Events & Campaign Tagging",
        "6. Enterprise Features (if applicable)",
        "7. Billing & Plan Upgrade Flow",
        "8. Admin Dashboard (Super Admin)"
    ]
    create_content_slide(prs, "Demo Walkthrough Plan", demo_steps)

    # Slide 15: Testing & Quality Assurance
    testing = [
        "Unit Tests: Jest + Supertest for API endpoints",
        "Integration Tests: Full workflow testing (scan → save → CRM)",
        "AI Pipeline Testing: Fallback mechanism validation",
        "UI Testing: Component rendering and user interaction",
        "Performance Testing: Batch processing and concurrent users",
        "Security Testing: Authentication, authorization, input validation"
    ]
    create_content_slide(prs, "Testing & Quality Assurance", testing)

    # Slide 16: Challenges & Solutions
    challenges = [
        "AI Accuracy: Implemented multi-level fallback (Gemini → OpenAI → Tesseract)",
        "Real-time Processing: Socket.IO for progress updates and collaboration",
        "Data Quality: Deduplication algorithms and manual merge workflows",
        "Scalability: Queued processing for batch operations",
        "Security: JWT authentication, input sanitization, audit logging",
        "Offline Capability: Tesseract.js fallback for single-card processing"
    ]
    create_content_slide(prs, "Challenges Faced & Solutions", challenges)

    # Slide 17: Future Enhancements
    future_features = [
        "Mobile App (React Native)",
        "Advanced AI: Custom model training, better accuracy",
        "Integration APIs: Salesforce, HubSpot, LinkedIn",
        "Analytics Dashboard: Lead conversion tracking, ROI metrics",
        "Multi-language Support: OCR for international business cards",
        "Advanced Compliance: GDPR, CCPA automation tools",
        "Workflow Automation: Zapier-style integration builder"
    ]
    create_content_slide(prs, "Future Scope & Enhancements", future_features)

    # Slide 18: Conclusion
    conclusion = [
        "Successfully built a comprehensive AI-powered CRM platform",
        "End-to-end implementation from concept to working product",
        "Modern tech stack with robust architecture and testing",
        "Enterprise-ready features with proper security and compliance",
        "Scalable foundation for future enhancements and integrations",
        "Real-world solution addressing actual business pain points"
    ]
    create_content_slide(prs, "Project Conclusion", conclusion)

    # Slide 19: Q&A
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    title_box = slide.shapes.title
    content_box = slide.placeholders[1]

    title_box.text = "Questions & Answers"
    title_box.text_frame.paragraphs[0].font.size = Pt(44)
    title_box.text_frame.paragraphs[0].font.bold = True

    qa_text = [
        "Thank you for your attention!",
        "",
        "We welcome your questions and feedback.",
        "",
        "Project Repository: Available for demonstration",
        "Live Demo: Ready for walkthrough",
        "Documentation: Complete technical documentation prepared"
    ]

    for item in qa_text:
        p = content_box.text_frame.add_paragraph()
        p.text = item
        p.font.size = Pt(24)
        if "Thank you" in item or "Questions" in item:
            p.font.bold = True

    return prs

def main():
    """Main function to generate the presentation"""
    print("Creating IntelliScan presentation...")

    # Create the presentation
    prs = create_intelliscan_presentation()

    # Save the presentation
    output_path = "IntelliScan_Presentation.pptx"
    prs.save(output_path)

    print(f"Presentation created successfully: {output_path}")
    print(f"Total slides: {len(prs.slides)}")

    # Print slide titles for verification
    print("\nSlide titles:")
    for i, slide in enumerate(prs.slides, 1):
        if slide.shapes.title:
            print(f"{i}. {slide.shapes.title.text}")
        else:
            print(f"{i}. [No title]")

if __name__ == "__main__":
    main()