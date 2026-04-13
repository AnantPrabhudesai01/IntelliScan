from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
from pathlib import Path

from PIL import Image


def build_deck(out_path: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Theme
    NAVY = RGBColor(8, 23, 48)
    WHITE = RGBColor(255, 255, 255)
    MUTED = RGBColor(200, 210, 225)
    ACCENT = RGBColor(64, 196, 255)
    ACCENT2 = RGBColor(110, 231, 183)

    TITLE_FONT = "Calibri"
    BODY_FONT = "Calibri"
    ASSETS_DIR = Path(__file__).resolve().parent / "extracted" / "intelliscan-pdf"

    def set_bg(slide, color=NAVY):
        bg = slide.background
        fill = bg.fill
        fill.solid()
        fill.fore_color.rgb = color

    def add_top_bar(slide):
        bar = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Pt(10)
        )
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()

    def add_title(slide, title, subtitle=None):
        tb = slide.shapes.add_textbox(Inches(0.9), Inches(1.7), Inches(11.6), Inches(2.2))
        tf = tb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = TITLE_FONT
        r.font.size = Pt(54)
        r.font.bold = True
        r.font.color.rgb = WHITE

        if subtitle:
            p2 = tf.add_paragraph()
            p2.space_before = Pt(14)
            r2 = p2.add_run()
            r2.text = subtitle
            r2.font.name = BODY_FONT
            r2.font.size = Pt(20)
            r2.font.color.rgb = MUTED

    def add_image_slide(slide, title: str, image_path: Path):
        # Title
        t = slide.shapes.add_textbox(Inches(0.9), Inches(0.65), Inches(12), Inches(0.8))
        tf = t.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = TITLE_FONT
        r.font.size = Pt(34)
        r.font.bold = True
        r.font.color.rgb = WHITE

        # Image fit box
        left = Inches(0.9)
        top = Inches(1.55)
        max_w = Inches(12.0)
        max_h = Inches(5.75)

        if not image_path.exists():
            # fallback text if missing
            missing = slide.shapes.add_textbox(left, top, max_w, Inches(1.0))
            m_tf = missing.text_frame
            m_tf.text = f"Missing image: {image_path.name}"
            m_tf.paragraphs[0].font.name = BODY_FONT
            m_tf.paragraphs[0].font.size = Pt(18)
            m_tf.paragraphs[0].font.color.rgb = MUTED
            return

        with Image.open(image_path) as im:
            w_px, h_px = im.size
        img_ratio = w_px / h_px if h_px else 1.0
        box_ratio = float(max_w) / float(max_h)

        if img_ratio >= box_ratio:
            pic_w = max_w
            pic_h = int(float(max_w) / img_ratio)
        else:
            pic_h = max_h
            pic_w = int(float(max_h) * img_ratio)

        # Center in box
        x = int(left + (max_w - pic_w) / 2)
        y = int(top + (max_h - pic_h) / 2)

        slide.shapes.add_picture(str(image_path), x, y, width=pic_w, height=pic_h)

    def add_section(slide, label, title):
        lb = slide.shapes.add_textbox(Inches(0.9), Inches(1.4), Inches(12), Inches(0.6))
        tf = lb.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = label
        r.font.name = BODY_FONT
        r.font.size = Pt(20)
        r.font.bold = True
        r.font.color.rgb = ACCENT

        tb = slide.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(12), Inches(2.5))
        tf2 = tb.text_frame
        tf2.clear()
        p2 = tf2.paragraphs[0]
        r2 = p2.add_run()
        r2.text = title
        r2.font.name = TITLE_FONT
        r2.font.size = Pt(64)
        r2.font.bold = True
        r2.font.color.rgb = WHITE

    def add_bullets(slide, title, bullets, x=Inches(0.9), y=Inches(1.2), w=Inches(12), h=Inches(5.8)):
        t = slide.shapes.add_textbox(x, y, w, Inches(0.8))
        tf = t.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = TITLE_FONT
        r.font.size = Pt(36)
        r.font.bold = True
        r.font.color.rgb = WHITE

        b = slide.shapes.add_textbox(x, y + Inches(0.95), w, h - Inches(0.95))
        tfb = b.text_frame
        tfb.word_wrap = True
        tfb.clear()

        for i, item in enumerate(bullets):
            if isinstance(item, (list, tuple)):
                parent = item[0]
                children = list(item[1:])
                p = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
                p.text = parent
                p.level = 0
                p.font.name = BODY_FONT
                p.font.size = Pt(24)
                p.font.color.rgb = MUTED
                for c in children:
                    pc = tfb.add_paragraph()
                    pc.text = c
                    pc.level = 1
                    pc.font.name = BODY_FONT
                    pc.font.size = Pt(20)
                    pc.font.color.rgb = MUTED
            else:
                p = tfb.paragraphs[0] if i == 0 else tfb.add_paragraph()
                p.text = item
                p.level = 0
                p.font.name = BODY_FONT
                p.font.size = Pt(24)
                p.font.color.rgb = MUTED

    def add_table(slide, title, headers, rows):
        t = slide.shapes.add_textbox(Inches(0.9), Inches(0.8), Inches(12), Inches(0.8))
        tf = t.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = TITLE_FONT
        r.font.size = Pt(36)
        r.font.bold = True
        r.font.color.rgb = WHITE

        cols = len(headers)
        table_shape = slide.shapes.add_table(
            len(rows) + 1, cols, Inches(0.9), Inches(1.8), Inches(12.0), Inches(5.2)
        )
        table = table_shape.table

        header_fill = RGBColor(20, 46, 86)
        cell_fill = RGBColor(12, 33, 66)

        for j, htxt in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = htxt
            for para in cell.text_frame.paragraphs:
                for run in para.runs:
                    run.font.name = BODY_FONT
                    run.font.size = Pt(16)
                    run.font.bold = True
                    run.font.color.rgb = WHITE
            cell.fill.solid()
            cell.fill.fore_color.rgb = header_fill

        for i, row in enumerate(rows, start=1):
            for j, val in enumerate(row):
                cell = table.cell(i, j)
                cell.text = val
                for para in cell.text_frame.paragraphs:
                    for run in para.runs:
                        run.font.name = BODY_FONT
                        run.font.size = Pt(16)
                        run.font.color.rgb = WHITE
                cell.fill.solid()
                cell.fill.fore_color.rgb = cell_fill

    def add_workflow(slide, title, steps):
        t = slide.shapes.add_textbox(Inches(0.9), Inches(0.8), Inches(12), Inches(0.8))
        tf = t.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = title
        r.font.name = TITLE_FONT
        r.font.size = Pt(36)
        r.font.bold = True
        r.font.color.rgb = WHITE

        left = Inches(0.9)
        top = Inches(2.0)
        box_w = Inches(3.0)
        box_h = Inches(1.15)
        gap_x = Inches(0.35)
        gap_y = Inches(0.35)

        for idx, (hdr, desc) in enumerate(steps):
            col = idx % 4
            row = idx // 4
            x = left + col * (box_w + gap_x)
            y = top + row * (box_h + gap_y)
            box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, box_w, box_h)
            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(12, 33, 66)
            box.line.color.rgb = ACCENT
            tfb = box.text_frame
            tfb.clear()

            p1 = tfb.paragraphs[0]
            r1 = p1.add_run()
            r1.text = hdr
            r1.font.name = BODY_FONT
            r1.font.size = Pt(16)
            r1.font.bold = True
            r1.font.color.rgb = ACCENT2

            p2 = tfb.add_paragraph()
            r2 = p2.add_run()
            r2.text = desc
            r2.font.name = BODY_FONT
            r2.font.size = Pt(14)
            r2.font.color.rgb = WHITE

    def new_slide():
        s = prs.slides.add_slide(prs.slide_layouts[6])
        set_bg(s)
        add_top_bar(s)
        return s

    # Slide 1: Title
    s = new_slide()
    add_title(s, "IntelliScan", "AI-Powered Business Card Scanner & Intelligent CRM Automation")
    sub = s.shapes.add_textbox(Inches(0.95), Inches(4.45), Inches(12), Inches(1.4))
    tf = sub.text_frame
    p = tf.paragraphs[0]
    p.text = "Submitted by: Anant Prabhudesai (245300694048)"
    p.font.name = BODY_FONT
    p.font.size = Pt(18)
    p.font.color.rgb = MUTED
    p2 = tf.add_paragraph()
    p2.text = "Internal Guide: Prof. Nil Gosai | SVIT, Vasad | April 2026"
    p2.font.name = BODY_FONT
    p2.font.size = Pt(18)
    p2.font.color.rgb = MUTED

    # Slide 2: Agenda
    s = new_slide()
    add_bullets(
        s,
        "Agenda",
        [
            "01  Introduction and Industry Context",
            "02  Overview of the Platform",
            "03  Core Platform Features",
            "04  Workflow of the System",
            "05  Tools & Technology Used",
            "06  System Design",
            "07  Development",
            "08  Agile Documentation",
            "09  Proposed Enhancements",
            "10  Conclusion",
        ],
        y=Inches(0.9),
    )

    # Intro section + content
    s = new_slide()
    add_section(s, "Introduction & Industry", "Context")

    s = new_slide()
    add_bullets(
        s,
        "Introduction & Industry Context",
        [
            "B2B networking still begins in the physical world (events, meetings, conferences).",
            "Manual CRM entry causes delays, errors, and lost opportunities.",
            "Sales teams need instant capture + deduplication + routing + follow-ups in one workflow.",
            "LLM + Vision OCR enables accurate extraction and demographic inference at scan-time.",
        ],
    )

    s = new_slide()
    add_bullets(
        s,
        "About IntelliScan (Overview)",
        [
            "Enterprise-grade web platform for end-to-end lead capture and pipeline execution.",
            "AI card scanning with OCR + Gemini Vision to extract text and infer missing demographics.",
            "Shared workspace (Rolodex + Kanban pipeline) to prevent double-outreach.",
            "Email marketing hub for AI-authored follow-ups and scheduled drip sequences.",
            "Public networking tools (public profile + calendar links) to convert leads into meetings.",
        ],
    )

    # Core Features section + slide
    s = new_slide()
    add_section(s, "Core Platform", "Features")

    s = new_slide()
    add_bullets(
        s,
        "Core Platform Features",
        [
            (
                "AI Card Scanning & Deduction Engine",
                "Multimodal extraction (image → structured contact)",
                "Industry/Seniority inference to enrich lead profiles",
            ),
            (
                "Shared Rolodex & CRM Pipeline",
                "Workspace-wide contact store with deduplication alerts",
                "Kanban-style pipeline tracking and stage progression",
            ),
            (
                "Email Marketing Hub",
                "AI Composer for personalized follow-ups",
                "Drip sequences + engagement telemetry (opens/clicks)",
            ),
            (
                "Public Networking Tools",
                "Public profile (digital business card)",
                "Calendar booking links for frictionless scheduling",
            ),
        ],
    )

    # Workflow slide
    s = new_slide()
    add_workflow(
        s,
        "End-to-End Workflow",
        [
            ("Scan / Upload", "Camera or file upload"),
            ("AI Extract", "OCR + structured fields"),
            ("Infer", "Industry & seniority"),
            ("Deduplicate", "Workspace match alerts"),
            ("Route", "Assign to rep / pipeline"),
            ("Manage", "Kanban stages & notes"),
            ("Compose", "AI email drafting"),
            ("Schedule", "Automated sequences"),
        ],
    )

    # Tools & Tech table
    s = new_slide()
    add_table(
        s,
        "Tools & Technologies Used",
        ["Frontend", "Backend", "Database", "AI Integration", "Dev & Testing"],
        [
            ["React (Vite)", "Node.js + Express", "SQLite", "Google Gemini Pro Vision", "Git + Postman"],
            ["Tailwind CSS", "REST APIs", "Relational schema", "OCR + demographic inference", "npm"],
        ],
    )

    # System Design section + key slides
    s = new_slide()
    add_section(s, "System", "Design")

    s = new_slide()
    add_bullets(
        s,
        "Target Users & Roles",
        [
            "Enterprise Sales Teams",
            "Event Marketing / HR Recruitment Teams",
            "Executive Networking Professionals",
            "RBAC roles: Super Admin, Enterprise Admin, Workspace Member, Anonymous (Public Profile)",
        ],
    )

    # Use case diagrams (from PDF)
    s = new_slide()
    add_image_slide(s, "Use Case Diagram (Entire System + Sales)", ASSETS_DIR / "p017_render.png")

    s = new_slide()
    add_image_slide(s, "Use Case Diagram (Enterprise Admin + Super Admin)", ASSETS_DIR / "p018_render.png")

    s = new_slide()
    add_image_slide(s, "Activity Diagram (Public Booking / Calendar Flow)", ASSETS_DIR / "p021_render.png")

    s = new_slide()
    add_image_slide(s, "Sequence Diagram (Business Card AI Vision Pipeline)", ASSETS_DIR / "p026_render.png")

    s = new_slide()
    add_image_slide(s, "Class Diagram (Authentication & Context Pipeline)", ASSETS_DIR / "p029_render.png")

    # Data dictionary
    s = new_slide()
    add_table(
        s,
        "Data Dictionary (Key Tables)",
        ["Table", "Primary Purpose", "Key Fields", "Constraints", "Notes"],
        [
            ["Workspaces", "Tenant/company context", "workspace_id, company_name", "PK, NN", "quota_scans, max_users"],
            ["Users", "Authentication + RBAC", "user_id, workspace_id, email", "PK, FK, UQ", "password_hash, is_active"],
        ],
    )

    s = new_slide()
    add_table(
        s,
        "Data Dictionary (Leads & Scans)",
        ["Table", "Primary Purpose", "Key Fields", "Constraints", "Notes"],
        [
            ["Contacts", "Lead record", "contact_id, workspace_id", "PK, FK", "pipeline_stage, confidence_score"],
            ["Scanned_Images", "Raw image & OCR", "image_id, contact_id", "PK, FK", "raw_ocr_json"],
        ],
    )

    s = new_slide()
    add_table(
        s,
        "Data Dictionary (Automation)",
        ["Table", "Primary Purpose", "Key Fields", "Constraints", "Notes"],
        [
            ["Events", "Campaign grouping", "event_id, workspace_id", "PK, FK", "custom_qr_url"],
            ["Email_Campaign_Nodes", "Email queue", "campaign_id, target_contact_id", "PK, FK", "delivery_status"],
            ["Calendar_Bookings", "Public scheduling", "booking_id, host_user_id", "PK, FK", "start/end blocks"],
        ],
    )

    # Development / UI section
    s = new_slide()
    add_section(s, "Development", "(UI Screens)")

    # UI screenshots rendered from the documentation PDF
    for title, img in [
        ("UI Screenshot: Landing / Login Screen", "p048_render.png"),
        ("UI Screenshot: Login + Sign Up", "p049_render.png"),
        ("UI Screenshot: Scanning Page + Leaderboard", "p050_render.png"),
        ("UI Screenshot: AI Coach + System Super Admin", "p051_render.png"),
        ("UI Screenshot: Organization + System Health", "p052_render.png"),
        ("UI Screenshot: Incident Center + Feedback", "p053_render.png"),
        ("UI Screenshot: Engine Performance", "p054_render.png"),
        ("UI Screenshot: Integrations + Versioning", "p055_render.png"),
        ("UI Screenshot: Calendar", "p056_render.png"),
    ]:
        s = new_slide()
        add_image_slide(s, title, ASSETS_DIR / img)

    # Agile section
    s = new_slide()
    add_section(s, "Agile", "Methodology")

    s = new_slide()
    add_bullets(
        s,
        "Agile Project Charter (Summary)",
        [
            "Project Title: IntelliScan",
            "Duration: 02/01/2026 – 30/04/2026",
            "Deliverables: AI Scanner, CRM Pipeline, AI Email Sequences, Event/Kiosk Module",
            "Risks: API rate limits/latency, complex React state, low-quality images",
            "Assumptions: stable internet, SMTP credentials available, adequate server resources",
        ],
    )

    s = new_slide()
    add_bullets(
        s,
        "Agile Roadmap",
        [
            "Q1 (02/01–16/01): Requirements + architecture + React/Node setup",
            "Q2 (19/01–30/01): DB schema + JWT auth + multi-tenancy",
            "Q3 (02/02–13/02): Gemini Vision integration (single & batch scanning)",
            "Q4 (16/02–27/02): CRM pipeline + semantic search + export",
            "Q5 (02/03–13/03): Email marketing hub + automation engine",
            "Q6 (16/03–27/03): Calendar + kiosk mode + enterprise theme",
            "Q7 (30/03–10/04): Testing + performance + deployment configs",
            "Q8 (13/04–30/04): Final documentation + submission",
        ],
    )

    s = new_slide()
    add_bullets(
        s,
        "Agile Project Plan",
        [
            "Release: v1.0",
            "Goals: architect an enterprise-grade AI SaaS and optimize heavy data workflows",
            "Core tasks: React state management, Gemini APIs, SQLite backend, integration testing, UI/UX rollout",
            "Success criteria: all modules integrated; documentation clarity; OCR accuracy",
        ],
    )

    s = new_slide()
    add_bullets(
        s,
        "User Stories (Core Modules)",
        [
            "AI Scanning Module: image → mapped CRM fields automatically",
            "CRM Data Pipeline: repository + smart semantic search + exports (CSV/vCard)",
            "Automated AI Sequences: personalized follow-ups based on inferred context",
            "Event & Kiosk Module: bulk lead capture during events + assignment to campaigns",
        ],
    )

    s = new_slide()
    add_bullets(
        s,
        "Sprint Backlog (Completed Highlights)",
        [
            "Sprint 1: requirement analysis, React setup, DB schema + UI prototypes",
            "Sprint 2: JWT auth + multi-tenancy; backend routing logic",
            "Sprint 3: Gemini Vision integration + scanning endpoint UI wiring",
            "Sprint 4: CRM pipeline + filtering; semantic search",
            "Sprint 5: AI email sequences; responsive CRM tables",
            "Sprint 6: calendar + kiosk mode; enterprise UI theme rollout",
            "Sprint 7: integration/system testing; academic documentation",
        ],
    )

    s = new_slide()
    add_table(
        s,
        "Earned Value Summary",
        ["Sprint", "Planned Hours", "Actual Hours", "Earned Value"],
        [
            ["Sprint 1", "240", "240", "10"],
            ["Sprint 2", "240", "200", "8.3"],
            ["Sprint 3", "240", "200", "8.3"],
            ["Sprint 4", "240", "336", "14"],
            ["Sprint 5", "240", "240", "10"],
            ["Sprint 6", "240", "190", "7.9"],
            ["Sprint 7", "120", "100", "8.3"],
        ],
    )

    # Enhancements + conclusion
    s = new_slide()
    add_bullets(
        s,
        "Proposed Enhancements",
        [
            "Secure connectors to external CRMs (Salesforce/HubSpot/SAP) with scoped permissions",
            "Advanced routing rules and scoring (industry + seniority + engagement)",
            "Improved batch processing and image-quality feedback",
            "Workspace collaboration improvements (audit logs, mentions, shared notes)",
            "Scaling to very large datasets (indexing + background jobs)",
        ],
    )

    s = new_slide()
    add_bullets(
        s,
        "Conclusion",
        [
            "IntelliScan bridges physical networking and digital sales execution.",
            "Unifies scanning, CRM pipeline, and automated follow-ups in one workspace.",
            "Reduces manual effort, speeds outreach, and prevents duplicate contact collisions.",
            "Built with modern full-stack tools and Gemini Vision-powered intelligence extraction.",
        ],
    )

    s = new_slide()
    add_title(s, "Thank You", "Questions?")

    prs.save(out_path)


if __name__ == "__main__":
    build_deck("IntelliScan_Final_Presentation.pptx")
